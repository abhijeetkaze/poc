# This script extracts the first slide from Presentation.pptx and saves it as an image in the output folder.
# Requirements: python-pptx, Pillow, and optionally, cairosvg for SVG images.

import os
from pptx import Presentation
from PIL import Image, ImageDraw

PPTX_PATH = os.path.join(os.path.dirname(__file__), 'resources', 'Presentation.pptx')
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), 'output')
OUTPUT_IMAGE = os.path.join(OUTPUT_DIR, 'slide1.png')


def render_slide_as_image(slide, slide_width, slide_height, width=1280, height=720):
    """
    Render a slide as an image (very basic: only background color and text).
    This does NOT render images, charts, or complex shapes.
    """
    img = Image.new('RGB', (width, height), color='white')
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        left = int(shape.left * width / slide_width)
        top = int(shape.top * height / slide_height)
        right = int((shape.left + shape.width) * width / slide_width)
        bottom = int((shape.top + shape.height) * height / slide_height)

        # Draw text
        if shape.has_text_frame:
            text = shape.text_frame.text
            draw.text((left, top), text, fill='black')

        # Draw pictures (images)
        if shape.shape_type == 13:  # PICTURE
            image = shape.image
            image_bytes = image.blob
            from io import BytesIO
            try:
                with Image.open(BytesIO(image_bytes)) as pic:
                    pic = pic.convert('RGBA')
                    pic = pic.resize((right - left, bottom - top))
                    # Handle transparency
                    if pic.mode in ('RGBA', 'LA'):
                        # Create a mask for transparency
                        mask = pic.split()[-1]
                        img.paste(pic, (left, top), mask)
                    else:
                        img.paste(pic, (left, top))
            except Exception as e:
                print(f"Error rendering image: {e}")

        # Draw rectangles and ellipses (basic shapes)
        if shape.shape_type == 1:  # RECTANGLE
            draw.rectangle([left, top, right, bottom], outline='blue', width=2)
        if shape.shape_type == 9:  # ELLIPSE
            draw.ellipse([left, top, right, bottom], outline='green', width=2)
    return img


def main():
    prs = Presentation(PPTX_PATH)
    first_slide = prs.slides[6]
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    img = render_slide_as_image(first_slide, slide_width, slide_height)
    img.save(OUTPUT_IMAGE)
    print(f"First slide saved as image: {OUTPUT_IMAGE}")

if __name__ == "__main__":
    main()
